import pip
try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    from docx import Document
    from PySimpleGUI import PySimpleGUI as sg
except:
    pip.main(["install", "--user", "python-docx"])
    pip.main(["install", "--user", "python-pptx"])
    pip.main(["install", "--user", "PySimpleGUI"])
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    from docx import Document
    from PySimpleGUI import PySimpleGUI as sg
def hex_to_rgb(hex_value):
    assert len(hex_value) == 6
    return tuple(int(hex_value[i:i+2], 16) for i in (0, 2, 4))
def make_new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    content = slide.placeholders[0] if len(slide.placeholders) > 0 else None
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return content
def set_content(content, prs):
    #content.left = Inches(0.5)
    #content.top = Inches(0.5)
    content.width = prs.slide_width #- Inches(1)
    content.height = prs.slide_height #- Inches(1)
    return content
def align_text_left(powerpoint_file):
    prs = Presentation(powerpoint_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    prs.save(powerpoint_file)
def convert_word_to_powerpoint(word_file, output_powerpoint):
    doc = Document(word_file)
    prs = Presentation()
    for paragraph in doc.paragraphs:
        content = make_new_slide(prs)
        if content:
            content = set_content(content, prs)
            p = content.text_frame.paragraphs[0]
            texts = []
            if len(paragraph.text) > 215:
                texta = str(paragraph.text).replace("\xa0", ' ')
                for i in range(0, len(paragraph.text)+1, 215):
                    max_index = -1
                    textb = texta[:215]
                    for j in ['.', '!', '?']:
                        if j in textb:
                            max_index = max(max_index, textb.rfind(j))
                    if max_index == -1:
                        max_index = textb.rfind(' ')
                        x = 1
                    texta = texta[max_index+1:]
                    texts.append(textb[:max_index+1])
                if i < len(paragraph.text):
                    max_index = -1
                    textb = texta[:180]
                    for j in ['.', '!', '?']:
                        if j in textb:
                            max_index = max(max_index, textb.rfind(j))
                    if max_index == -1:
                        max_index = textb.rfind(' ')
                    texts.append(textb[:max_index+1])
                textlist = []
                for j in texts:
                    if j[0] == ' ' if len(j) > 0 else False:
                        j = j[1:]
                    textlist.append(j)
                texts = textlist
            for run in paragraph.runs:
                r = p.add_run()
                if (((len(content.text_frame.text) >= len(texts[0])-1) if texts != [] else False) and ('.' in run.text or '!' in run.text or '?' in run.text)) or (((len(content.text_frame.text) >= len(texts[0])-1) if texts != [] else False) and texts[0][-1] == ' ' and ' ' in run.text):
                    texts.pop(0)
                    if '.' in run.text:
                        r.text = '.'
                    elif '!' in run.text:
                        r.text = '!'
                    elif '?' in run.text:
                        r.text = '?'
                    content = make_new_slide(prs)
                    if content:
                        content = set_content(content, prs)
                        p = content.text_frame.paragraphs[0]
                    k = p.add_run()
                    if ('\u00A0' == run.text[1] or ' ' == run.text[1]) if len(run.text) > 1 else False:
                        k.text = run.text.replace('.', '').replace('!', '').replace('?', '')[1:]
                    if run.font.italic == True:
                        k.font.name = "GungsuhChe"
                    else:
                        k.font.name = "Malgun Gothic"
                    if hex_to_rgb(str(run.font.color.rgb)) == (255, 255, 0) or run.font.italic == True:
                        k.font.color.rgb = RGBColor(255, 255, 0)
                    else:
                        k.font.color.rgb = RGBColor(255, 255, 255)
                    k.font.size = Pt(44)
                    k.font.bold = True #run.font.bold
                    k.font.underline = run.font.underline
                else:
                    if len(content.text_frame.text) == 0:
                        if all(char == ' ' for char in run.text) or all(char == '\u00A0' for char in run.text):
                            r.text = ''
                        else:
                            r.text = run.text
                    else:
                        r.text = run.text
                if run.font.italic == True:
                    r.font.name = "GungsuhChe"
                else:
                    r.font.name = "Malgun Gothic"
                if hex_to_rgb(str(run.font.color.rgb)) == (255, 255, 0) or run.font.italic == True:
                    r.font.color.rgb = RGBColor(255, 255, 0)
                else:
                    r.font.color.rgb = RGBColor(255, 255, 255)
                r.font.size = Pt(44)
                r.font.bold = True #run.font.bold
                r.font.underline = run.font.underline
    prs.save(output_powerpoint)
def delete_empty_textbox_slides(powerpoint_file):
    prs = Presentation(powerpoint_file)
    slides_to_delete = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and (len(shape.text_frame.text) == 0 or len(shape.text_frame.text) == 1):
                slides_to_delete.append(i)
    slides_to_delete.reverse()
    for i in slides_to_delete:
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    prs.save(powerpoint_file)
sg.theme('DarkBlue3')
layout = [[sg.Text('Select Input File:'), sg.Input(key='_IN_FILE_'), sg.FileBrowse(file_types=(("Word files", "*.docx"), ("Word files", "*.doc")))],
          [sg.Text('Select Output File:'), sg.Input(key='_OUT_FILE_'), sg.FileSaveAs(file_types=(("PowerPoint files", "*.pptx"),), default_extension='.pptx')],
          [sg.Button('Convert'), sg.Button('Cancel')]]
window = sg.Window('DocToPPT Converter', layout)
while True:
    event, values = window.read()
    if event == sg.WIN_CLOSED or event == 'Cancel':
        break
    elif event == 'Convert':
        try:
            input_file = values['_IN_FILE_']
            output_file = values['_OUT_FILE_']
            convert_word_to_powerpoint(input_file, output_file)
            delete_empty_textbox_slides(output_file)
            align_text_left(output_file)
            sg.popup('Conversion Complete!')
            break
        except:
            sg.popup('Error: Please make sure that the input file is a Word file and that the output file is a PowerPoint file.')
window.close()